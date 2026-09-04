#!/usr/bin/env python3
"""Build the editable 'Alternative Approach' summary deck from the live data JSON.
All content is native PowerPoint shapes/text/tables — fully editable, no images."""
import json, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))
D = json.load(open(os.path.join(HERE, "deck-data.json")))
OUT = os.path.abspath(os.path.join(HERE, "..", "docs", "Alternative Approach - Summary.pptx"))

INK=RGBColor(0x16,0x18,0x1D); MUT=RGBColor(0x5C,0x60,0x66); ACC=RGBColor(0xC4,0x3B,0x44)
TEC=RGBColor(0x35,0x59,0x7B); GRN=RGBColor(0x2E,0x7D,0x4F); OPS=RGBColor(0xA3,0x67,0x1F)
GOV=RGBColor(0x68,0x4E,0x86); PEO=RGBColor(0x55,0x76,0x3F)
LINE=RGBColor(0xE5,0xE2,0xDB); PAPER=RGBColor(0xFB,0xFA,0xF8); WHITE=RGBColor(0xFF,0xFF,0xFF)
FONT="Aptos"; FALL="Calibri"

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height

def slide(bg=PAPER):
    s=prs.slides.add_slide(BLANK)
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,SW,SH)
    r.fill.solid(); r.fill.fore_color.rgb=bg; r.line.fill.background(); r.shadow.inherit=False
    return s

def txt(s,l,t,w,h,runs,size=14,color=INK,bold=False,align=PP_ALIGN.LEFT,sp=1.06,anchor=None):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=tb.text_frame; tf.word_wrap=True
    if anchor: tf.vertical_anchor=anchor
    lines=runs if isinstance(runs,list) else [runs]
    for i,line in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align; p.line_spacing=sp
        if isinstance(line,tuple):
            text=line[0]; csz=line[1] if len(line)>1 else size
            ccol=line[2] if len(line)>2 else color; cbold=line[3] if len(line)>3 else bold
        else:
            text,csz,ccol,cbold=line,size,color,bold
        r=p.add_run(); r.text=text; f=r.font; f.name=FONT; f.size=Pt(csz); f.color.rgb=ccol; f.bold=cbold
    return tb

def box(s,l,t,w,h,fill=WHITE,line=LINE,lw=0.75,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.06):
    c=s.shapes.add_shape(shape,Inches(l),Inches(t),Inches(w),Inches(h))
    if fill is None: c.fill.background()
    else: c.fill.solid(); c.fill.fore_color.rgb=fill
    if line is None: c.line.fill.background()
    else: c.line.color.rgb=line; c.line.width=Pt(lw)
    c.shadow.inherit=False
    try:
        if shape==MSO_SHAPE.ROUNDED_RECTANGLE: c.adjustments[0]=radius
    except Exception: pass
    return c

def topbar(s,color=ACC):
    b=box(s,0.7,1.24,0.62,0.05,fill=color,line=None,shape=MSO_SHAPE.RECTANGLE)
    return b

def head(s,kicker,title,tcolor=INK,num=None):
    k = (f"{num}   " if num else "") + kicker.upper()
    txt(s,0.7,0.42,12,0.3,(k,11,MUT,True))
    txt(s,0.7,0.70,12.0,0.7,(title,25,tcolor,True))
    topbar(s)

def chip(s,l,t,text,fill,tcolor=WHITE,w=None):
    w = w or (0.14+0.085*len(text))
    c=box(s,l,t,w,0.28,fill=fill,line=None,radius=0.5)
    tf=c.text_frame; tf.word_wrap=False; tf.margin_top=Pt(1); tf.margin_bottom=Pt(1)
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; r=p.add_run(); r.text=text
    r.font.name=FONT; r.font.size=Pt(9.5); r.font.bold=True; r.font.color.rgb=tcolor
    return c

def bullets(s,l,t,w,h,items,size=12.5,color=INK,gap=1.3,marker="•  ",mcolor=None):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=tb.text_frame; tf.word_wrap=True
    for i,it in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.line_spacing=gap; p.space_after=Pt(2)
        rm=p.add_run(); rm.text=marker; rm.font.name=FONT; rm.font.size=Pt(size); rm.font.bold=True
        rm.font.color.rgb=mcolor or color
        r=p.add_run(); r.text=it; r.font.name=FONT; r.font.size=Pt(size); r.font.color.rgb=color
    return tb

def table(s,l,t,w,rows,colw,header_fill=None,body_size=11,head_size=10.5,rowh=0.42):
    nr=len(rows); nc=len(rows[0])
    gt=s.shapes.add_table(nr,nc,Inches(l),Inches(t),Inches(w),Inches(rowh*nr)).table
    gt.first_row=False; gt.horz_banding=False
    for ci,cw in enumerate(colw): gt.columns[ci].width=Inches(cw)
    for ri,row in enumerate(rows):
        gt.rows[ri].height=Inches(rowh)
        for ci,val in enumerate(row):
            cell=gt.cell(ri,ci); cell.margin_left=Inches(0.08); cell.margin_right=Inches(0.06)
            cell.margin_top=Inches(0.03); cell.margin_bottom=Inches(0.03)
            cell.vertical_anchor=MSO_ANCHOR.MIDDLE
            if ri==0:
                cell.fill.solid(); cell.fill.fore_color.rgb=PAPER
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb=WHITE
            tf=cell.text_frame; tf.word_wrap=True; p=tf.paragraphs[0]
            r=p.add_run(); r.text=str(val); f=r.font; f.name=FONT
            f.size=Pt(head_size if ri==0 else body_size)
            f.bold=(ri==0 or ci==0); f.color.rgb=(MUT if ri==0 else INK)
    return gt

def foot(s,text,color=GRN):
    txt(s,0.7,7.0,12,0.35,(text,12,color,True))

# ══════════════════════════════════════════════════════════════════════
# 1 · TITLE
s=slide()
box(s,0,0,SW,SH,fill=PAPER,line=None,shape=MSO_SHAPE.RECTANGLE)
txt(s,0.9,0.9,11.5,0.35,("PANDORA TS&F · STRATEGIC PARTNER SELECTION 2027 · THE REQUESTED ALTERNATIVE",12,ACC,True))
txt(s,0.9,1.5,11.4,1.5,[("The Alternative Approach",40,INK,True),("Run As-Is First, Earn Agentic",40,INK,True)],sp=1.02)
box(s,0.92,3.35,3.4,0.6,fill=None,line=ACC,lw=1.5)
txt(s,1.05,3.4,10.6,0.9,('"We found the current approach to be a bit ambitious. Can you propose an alternative?"',15,MUT,False),sp=1.15)
txt(s,0.9,4.5,11,0.35,("Same destination. A gentler on-ramp. The fabric is the destination — never a dependency.",15,TEC,True))
txt(s,0.9,5.15,11.4,1.0,("In the deep dive we showed the destination, not the road to it — so it read as if we would do it all from day one. This is the road: we run your platforms as they run today, and travel to the agentic destination step by earned step, at your pace, behind gates you control.",13,MUT),sp=1.2)
txt(s,0.9,6.7,11,0.3,("Anchored to contract award, 30 October 2026",11,MUT))

# 2 · WHAT CHANGES
s=slide(); head(s,"What we heard — and what changes","A real alternative, not repackaging",num="01")
rows=[["Dimension","July approach","This alternative"]]+[list(r) for r in D["changesTable"]]
table(s,0.7,1.55,12.0,rows,[3.4,4.3,4.3],rowh=0.62,body_size=11)
foot(s,"What does not change: the destination, the SLAs, the assets, the RFP commitments. The July approach remains available as the accelerated option.")

# 3 · TWO-LANE MODEL + GATES
s=slide(); head(s,"The model","Two lanes, one destination — three gate kinds",num="02")
box(s,0.7,1.65,12,0.95,fill=RGBColor(0xED,0xF1,0xF6),line=RGBColor(0xB6,0xC5,0xD5))
txt(s,0.95,1.74,11.5,0.4,("Lane 1 · Run & Deliver — active from day one",15,TEC,True))
txt(s,0.95,2.12,11.5,0.4,("Transition per the RFP · your tools & SLAs · migration factories · conventional automation · foundations built as a by-product",11.5,INK))
box(s,0.7,2.75,12,0.95,fill=RGBColor(0xF7,0xE9,0xE9),line=RGBColor(0xE4,0xB9,0xBD))
txt(s,0.95,2.84,11.5,0.4,("Lane 2 · Improve & Evolve — the agentic fabric, held ready",15,ACC,True))
txt(s,0.95,3.22,11.5,0.4,("Nothing in the estate until Gate 0 + the item gate pass · per-item dial in Pandora's hand · reversible, audited, quarterly-reviewed",11.5,INK))
box(s,0.7,3.85,12,0.5,fill=RGBColor(0xF1,0xED,0xF6),line=RGBColor(0xC8,0xB9,0xDA))
txt(s,0.95,3.92,11.5,0.4,("Governance base — service governance active day one · agentic governance pre-built, dormant until the first dial-up",11.5,GOV,True))
# gate series row
gx=0.7
for i,g in enumerate(D["gateSeries"]):
    x=0.7+i*4.05
    box(s,x,4.65,3.85,1.75,fill=WHITE,line=GRN,lw=1.0)
    txt(s,x+0.18,4.75,3.5,0.35,(g["name"],13,GRN,True))
    txt(s,x+0.18,5.12,3.5,0.3,(g["scope"],10.5,MUT,True))
    txt(s,x+0.18,5.42,3.5,0.9,(g["unlocks"]+".",11,INK),sp=1.1)
foot(s,"Only Gate 0 asks us to build something — so only Gate 0 gets its own slide (next). Gate 1 & Gate 2 are the short per-item checks above; every gate is Pandora-owned and reversible.")

# 4 · SCOPE
s=slide(); head(s,"The scope — what we are talking about","Your scope of work: 30 items, 13 DevOps + 17 Data & Integration",num="03")
txt(s,0.7,1.5,12.2,0.55,("Taken directly from your platform deep-dive scope mapping — nothing added, nothing dropped. We only make the service level of each item something you dial.",12.5,MUT),sp=1.15)
dev=[i for i in D["scopeDial"] if i["platform"]=="devops"]; dat=[i for i in D["scopeDial"] if i["platform"]=="data"]
box(s,0.7,2.2,6.0,4.55,line=LINE); box(s,0.7,2.2,6.0,0.05,fill=TEC,line=None,shape=MSO_SHAPE.RECTANGLE)
txt(s,0.9,2.32,5.6,0.3,("DEVOPS PLATFORM · 13 ITEMS",11,TEC,True))
txt(s,0.9,2.7,5.7,4.0,"\n".join(f"{i}.  {it['name']}" for i,it in enumerate(dev,1)),11,INK,sp=1.12)
box(s,6.95,2.2,5.75,4.55,line=LINE); box(s,6.95,2.2,5.75,0.05,fill=GOV,line=None,shape=MSO_SHAPE.RECTANGLE)
txt(s,7.15,2.32,5.4,0.3,("DATA & INTEGRATION PLATFORM · 17 ITEMS",11,GOV,True))
txt(s,7.15,2.7,5.5,4.0,"\n".join(f"{i}.  {it['name']}" for i,it in enumerate(dat,1)),10,INK,sp=1.03)

# 5 · FOUNDATIONS / GATE 0
s=slide(); head(s,"Gate 0 in depth","The one gate that needs building — four readiness layers that must hold",num="04")
txt(s,0.7,1.5,12.2,0.6,("You just met the three gates; this is the one that needs building — which is why it, and only it, gets a slide of its own. Lane 1 builds all four layers as a by-product of running the service well, so running as-is first is not a delay on the journey — it is its first leg.",12.5,MUT),sp=1.15)
# capstone
box(s,3.9,2.25,5.5,0.5,fill=GRN,line=None); txt(s,3.9,2.32,5.5,0.4,("GATE 0 · all four layers proven → unlocks Lane 2",12.5,WHITE,True),align=PP_ALIGN.CENTER)
hues=[GOV,OPS,PEO,TEC]
for i,pil in enumerate(D["foundationPillars"][::-1]):  # top-down: governance..infra
    y=2.95+i*0.82
    box(s,2.3,y,8.7,0.72,fill=WHITE,line=LINE); box(s,2.3,y,0.06,0.72,fill=hues[i],line=None,shape=MSO_SHAPE.RECTANGLE)
    bx=box(s,2.45,y+0.13,0.46,0.46,fill=hues[i],line=None); txt(s,2.45,y+0.19,0.46,0.35,(str(i+1),15,WHITE,True),align=PP_ALIGN.CENTER)
    txt(s,3.05,y+0.09,7.8,0.3,(pil["title"],12.5,INK,True))
    txt(s,3.05,y+0.36,7.8,0.3,(pil["why"],10.5,MUT,False))
box(s,2.3,6.3,8.7,0.42,fill=INK,line=None); txt(s,2.3,6.36,8.7,0.32,("Your estate, day one — Lane 1 builds all four layers as it runs",11,RGBColor(0xE8,0xE6,0xE1),True),align=PP_ALIGN.CENTER)

# 6 · DAY ONE TRANSITION
s=slide(); head(s,"Run as-is","Day one: the RFP delivered as written — anchored to the 30 Oct award",num="05")
tp=D["transitionPlan"]
for i,ph in enumerate(tp):
    y=1.65+i*0.82
    box(s,0.7,y,12,0.72)
    txt(s,0.9,y+0.13,2.7,0.5,(ph["phase"],12.5,INK,True))
    txt(s,3.7,y+0.13,2.6,0.5,(ph["window"],11.5,ACC,True))
    txt(s,6.4,y+0.11,6.1,0.55,(ph["detail"],10.5,MUT),sp=1.1)
txt(s,0.7,5.1,12.2,1.7,[(f["label"]+" — "+f["value"],11.5,INK,False) for f in D["dayOneFacts"]],sp=1.25)

# 7 · PEOPLE / TECH / OPS
s=slide(); head(s,"Run as-is","People · Technology · Operations — as-is first",num="06")
for i,ln in enumerate(D["asIsLanes"]):
    x=0.7+i*4.05
    box(s,x,1.65,3.85,4.6)
    txt(s,x+0.2,1.78,3.5,0.4,(ln["title"],14.5,TEC,True))
    bullets(s,x+0.2,2.25,3.5,3.0,ln["points"],size=10.5,gap=1.15)
    box(s,x+0.15,5.55,3.55,0.6,fill=RGBColor(0xE9,0xF3,0xEC),line=None); txt(s,x+0.28,5.62,3.35,0.5,(ln["message"],11,GRN,True),sp=1.05)

# 8 · GOVERNANCE
s=slide(); head(s,"Governance","Two layers, one principle — controls arrive before autonomy",num="07")
gl=D["governanceLayers"]
box(s,0.7,1.7,12,2.35,fill=RGBColor(0xE9,0xF3,0xEC),line=GRN)
txt(s,0.95,1.82,11.4,0.4,(gl[0]["name"]+" — "+gl[0]["state"],15,GRN,True))
txt(s,0.95,2.2,11.5,0.35,(gl[0]["sub"],11,MUT))
bullets(s,0.95,2.55,11.4,1.4,gl[0]["items"],size=11.5,gap=1.2,mcolor=GRN)
box(s,0.7,4.2,12,2.35,fill=PAPER,line=OPS)
txt(s,0.95,4.32,11.4,0.4,(gl[1]["name"]+" — "+gl[1]["state"],15,OPS,True))
txt(s,0.95,4.7,11.5,0.35,(gl[1]["sub"],11,MUT))
bullets(s,0.95,5.05,11.4,1.4,gl[1]["items"],size=11.5,gap=1.2,mcolor=OPS)
foot(s,"The quarterly dial review joins the two layers — every item's level, decided by Pandora. \"Not yet\" is always valid.")

# 9 · THE TEAM — LOCATIONS
s=slide(); head(s,"The Team","One team, three locations — not a hand-off",num="08")
txt(s,0.7,1.5,12.2,0.5,("One team · one backlog · one leader · one knowledge base",12.5,MUT,True))
for i,loc in enumerate(D["teamLocations"]):
    x=0.7+i*4.05
    box(s,x,2.1,3.85,4.3); box(s,x,2.1,3.85,0.05,fill=TEC,line=None,shape=MSO_SHAPE.RECTANGLE)
    txt(s,x+0.2,2.22,3.5,0.35,(loc["city"],16,INK,True))
    txt(s,x+0.2,2.6,3.5,0.3,(loc["kind"],10.5,MUT,True))
    txt(s,x+0.2,2.95,3.5,0.7,(loc["role"],11,INK),sp=1.1)
    bullets(s,x+0.2,3.7,3.5,1.5,loc["roles"],size=10.5,gap=1.15)
    l1=loc["lane1"]; l2=100-l1
    txt(s,x+0.2,5.35,3.5,0.22,("WHERE ITS EFFORT GOES TODAY",8,MUT,True))
    box(s,x+0.2,5.63,0.14,0.14,fill=TEC,line=None,shape=MSO_SHAPE.RECTANGLE); txt(s,x+0.42,5.59,3.2,0.22,(f"Lane 1 · Run & Deliver  {l1}%",9.5,INK,True))
    box(s,x+0.2,5.89,0.14,0.14,fill=ACC,line=None,shape=MSO_SHAPE.RECTANGLE); txt(s,x+0.42,5.85,3.2,0.22,(f"Lane 2 · Improve & Evolve  {l2}%",9.5,INK,True))
    box(s,x+0.2,6.14,3.5*l1/100.0,0.16,fill=TEC,line=None,shape=MSO_SHAPE.RECTANGLE)
    box(s,x+0.2+3.5*l1/100.0,6.14,max(0.12,3.5*l2/100.0),0.16,fill=ACC,line=None,shape=MSO_SHAPE.RECTANGLE)
foot(s,"Read every bar the same way: each location is mostly Lane 1 today. The coral Lane 2 slice is not a second team — it is the same people doing enablement, and it grows only as you turn dials.")

# 10 · ONE LEADER + CONTROL
s=slide(); head(s,"The Team","One leader across all tracks — and Pandora on the wheel",num="09")
box(s,3.2,1.55,7.0,0.7,fill=RGBColor(0xED,0xF1,0xF6),line=RGBColor(0xB6,0xC5,0xD5))
txt(s,3.2,1.62,7.0,0.3,("Pandora leadership — retains architecture · standards · roadmap · Lead & Senior Engineers",11.5,TEC,True),align=PP_ALIGN.CENTER)
txt(s,3.2,2.28,7.0,0.25,("↕ pairs 1:1",10,MUT,True),align=PP_ALIGN.CENTER)
box(s,3.9,2.55,5.6,0.62,fill=RGBColor(0xF7,0xE9,0xE9),line=RGBColor(0xE4,0xB9,0xBD))
txt(s,3.9,2.63,5.6,0.4,("One Sapient Delivery Lead — accountable across all tracks, both lanes",12,ACC,True),align=PP_ALIGN.CENTER)
for i,tr in enumerate(D["teamTracks"]):
    x=0.7+i*4.05
    box(s,x,3.4,3.85,1.5); txt(s,x+0.15,3.5,3.55,0.3,(tr["name"],12.5,INK,True),align=PP_ALIGN.CENTER)
    box(s,x+0.15,3.9,3.55,0.44,fill=RGBColor(0xEF,0xF3,0xEA),line=None); txt(s,x+0.25,3.96,3.4,0.35,("DEV  "+tr["dev"],9.5,PEO,False))
    box(s,x+0.15,4.4,3.55,0.44,fill=RGBColor(0xF8,0xF1,0xE6),line=None); txt(s,x+0.25,4.46,3.4,0.35,("OPS  "+tr["ops"],9.5,OPS,False))
# control strip
txt(s,0.7,5.15,12,0.3,("WHO DECIDES, WHO DOES — CONTROL BY DESIGN",11,MUT,True))
cb=D["controlBands"]; cw=12.0/len(cb); cmap={"pandora":TEC,"joint":GOV,"sapient":ACC}
for i,b in enumerate(cb):
    x=0.7+i*cw
    box(s,x+0.03,5.5,cw-0.06,0.6,fill=cmap[b["owner"]],line=None)
    txt(s,x+0.05,5.56,cw-0.1,0.5,(b["area"],9,WHITE,True),align=PP_ALIGN.CENTER,sp=1.0)
txt(s,0.7,6.25,12,0.3,("■ Pandora-owned    ■ Joint    ■ Sapient-run",10.5,MUT,True))
foot(s,"Control is structural: every gate and every dial-up needs a named Pandora owner's sign-off.")

# 11 · WHY TWO TEAMS — CONVERGENCE
s=slide(); head(s,"The Team","Why two teams? One team, sequenced — dev & ops converge",num="10")
txt(s,0.7,1.5,12.2,0.85,(D["convergeIntro"],11.5,MUT),sp=1.15)
# convergence: two boxes merging into one
box(s,0.7,2.5,2.4,0.6,fill=RGBColor(0xEF,0xF3,0xEA),line=PEO); txt(s,0.7,2.58,2.4,0.4,("Development",12,PEO,True),align=PP_ALIGN.CENTER)
box(s,0.7,3.5,2.4,0.6,fill=RGBColor(0xF8,0xF1,0xE6),line=OPS); txt(s,0.7,3.58,2.4,0.4,("Operations · 24×7",11,OPS,True),align=PP_ALIGN.CENTER)
a1=s.shapes.add_connector(2,Inches(3.1),Inches(2.8),Inches(9.0),Inches(3.3)); a1.line.color.rgb=PEO; a1.line.width=Pt(3)
a2=s.shapes.add_connector(2,Inches(3.1),Inches(3.8),Inches(9.0),Inches(3.35)); a2.line.color.rgb=OPS; a2.line.width=Pt(3)
box(s,9.1,2.95,3.5,0.85,fill=RGBColor(0xED,0xF1,0xF6),line=TEC); txt(s,9.1,3.02,3.5,0.7,[("One team",13,TEC,True),("build it · run it · agent-supervised",10,INK,False)],align=PP_ALIGN.CENTER,sp=1.0)
box(s,3.4,4.35,6.2,0.6,fill=RGBColor(0xE9,0xF3,0xEC),line=GRN); txt(s,3.45,4.42,6.1,0.5,("Agentic ops (Ops L2/L3) removes the operational cognitive load — so build-minded engineers can own what they run.",10.5,GRN,True),align=PP_ALIGN.CENTER,sp=1.0)
# mechanisms 2x2
for i,m in enumerate(D["convergeMechanisms"]):
    x=0.7+(i%2)*6.15; y=5.15+(i//2)*0.85
    box(s,x,y,5.95,0.78); txt(s,x+0.15,y+0.08,5.65,0.3,(m["title"],11,INK,True)); txt(s,x+0.15,y+0.34,5.65,0.4,(m["detail"],9.3,MUT),sp=1.05)

# 12 · SKILLS + KT
s=slide(); head(s,"The Team","Skills coverage — Ops L1→L3 into Development, and knowledge transfer that sticks",num="11")
# Full-width tier matrix: skill × [Ops L1, Ops L2, Ops L3, Development] — mirrors the site's skills ladder
thead=["Engineering skill"]+[f"{h['label']} · {h['detail']}" for h in D["skillTierHeaders"]]
rows=[thead]+[[r["skill"],r["t1"],r["t2"],r["t3"],r["dev"]] for r in D["skillRows"]]
table(s,0.7,1.55,12.0,rows,[2.2,2.45,2.45,2.45,2.45],rowh=0.5,body_size=7.8,head_size=8.5)
# KT loop as a horizontal strip beneath the matrix
txt(s,0.7,5.35,12.0,0.28,("KNOWLEDGE TRANSFER — A DESIGNED LOOP",10.5,GOV,True))
n=len(D["ktLoop"]); kw=(12.0-(n-1)*0.2)/n
for i,k in enumerate(D["ktLoop"]):
    x=0.7+i*(kw+0.2); c = GRN if k["step"]=="Own" else GOV
    box(s,x,5.7,kw,1.0,fill=(RGBColor(0xE9,0xF3,0xEC) if k["step"]=="Own" else WHITE),line=LINE)
    box(s,x,5.7,kw,0.05,fill=c,line=None,shape=MSO_SHAPE.RECTANGLE)
    txt(s,x+0.12,5.8,kw-0.24,0.3,(f"{i+1}. {k['step']}",10.5,c,True))
    txt(s,x+0.12,6.08,kw-0.24,0.55,(k["detail"],8.2,MUT),sp=1.0)
foot(s,"Full coverage day one across dev & ops — incl. the Kubernetes, Kafka & DevOps fundamentals your RFP flagged as thin.")

# 13 · CAPACITY
s=slide(); head(s,"The Team","Capacity that compounds — more from a leaner, stabler team",num="12")
# simple chart: axes + two lines via connectors
ox,oy,cw2,ch=1.2,2.2,7.5,3.0
box(s,ox,oy,cw2,ch,fill=WHITE,line=LINE)
cap=s.shapes.add_connector(2,Inches(ox+0.3),Inches(oy+2.5),Inches(ox+cw2-0.3),Inches(oy+0.5)); cap.line.color.rgb=TEC; cap.line.width=Pt(3)
hc=s.shapes.add_connector(2,Inches(ox+0.3),Inches(oy+2.0),Inches(ox+cw2-0.3),Inches(oy+2.15)); hc.line.color.rgb=MUT; hc.line.width=Pt(2.5); hc.line.dash_style=2
txt(s,ox+cw2-2.6,oy+0.25,2.4,0.3,("Effective capacity ↑",10.5,TEC,True))
txt(s,ox+0.4,oy+1.6,3.0,0.3,("Team size — flat / leaner",10.5,MUT,True))
txt(s,ox+3.0,oy+0.95,2.4,0.5,("productivity dividend",10,MUT,True))
for i,dt in enumerate(["Jan 2027","Apr 2027","Oct 2027","Oct 2028"]):
    x=ox+0.3+i*(cw2-0.6)/3.0; txt(s,x-0.6,oy+ch+0.05,1.2,0.3,(dt,9.5,INK,True),align=PP_ALIGN.CENTER)
# drivers on right
for i,dr in enumerate(D["capacityDrivers"]):
    y=2.2+i*0.78
    box(s,9.0,y,3.6,0.7); txt(s,9.15,y+0.08,3.35,0.3,(f"{i+1}. {dr['label']}  ({dr['at']})",10.5,TEC,True)); txt(s,9.15,y+0.33,3.35,0.35,(dr["detail"],8.8,MUT),sp=1.0)
foot(s,D["capacityNote"])

# 14 · THE DIAL — LEVELS & GATES
s=slide(); head(s,"The Dial","Three service levels, and the two gates between them",num="13")
box(s,0.7,1.6,12,0.55,fill=RGBColor(0xE9,0xF3,0xEC),line=GRN); txt(s,0.9,1.68,11.6,0.4,("GATE 0 · foundations proven — estate-wide, passes once. Only then can any per-item dial move.",12,GRN,True),align=PP_ALIGN.CENTER)
lv=D["dialLevels"]; lc=[TEC,OPS,ACC]
for i,l in enumerate(lv):
    x=0.7+i*4.05
    box(s,x,2.5,3.85,2.1); box(s,x,2.5,3.85,0.05,fill=lc[i],line=None,shape=MSO_SHAPE.RECTANGLE)
    txt(s,x+0.2,2.6,3.5,0.4,(f"L{l['level']} · {l['name']}",14,lc[i],True))
    txt(s,x+0.2,3.05,3.5,1.4,(l["detail"],10.5,MUT),sp=1.12)
    if i<2:
        g=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x+3.7),Inches(3.35),Inches(0.7),Inches(0.4)); g.fill.solid(); g.fill.fore_color.rgb=GRN; g.line.fill.background(); g.shadow.inherit=False
        txt(s,x+3.7,3.4,0.7,0.3,(f"G{i+1}",11,WHITE,True),align=PP_ALIGN.CENTER)
foot(s,"Every item starts at L0 and stays there until you turn its dial. You hold all 30 dials.")

# 15 · WALKTHROUGH
s=slide(); head(s,"The Dial","The dial in action — the <30-minute workspace, same service three settings",num="14")
wk=D["walkthroughs"][0]; shares=[15,55,85]; bl=["Human-run · engineer approves & applies","AI-assisted · human approves every action","Agent-run · human reviews outcomes"]
txt(s,0.7,1.55,12,0.4,(wk["title"],14,INK,True))
for i,st in enumerate(wk["steps"]):
    y=2.1+i*1.55
    box(s,0.7,y,12,1.4); box(s,0.7,y,0.06,1.4,fill=[TEC,OPS,ACC][i],line=None,shape=MSO_SHAPE.RECTANGLE)
    txt(s,0.95,y+0.12,4.5,0.3,(st["level"],12,[TEC,OPS,ACC][i],True))
    txt(s,5.6,y+0.12,7.0,0.3,(bl[i],10.5,MUT,True))
    # human/agent bar
    hw=11.7*(100-shares[i])/100.0
    box(s,0.95,y+0.5,hw,0.24,fill=TEC,line=None); box(s,0.95+hw,y+0.5,11.7*shares[i]/100.0,0.24,fill=ACC,line=None)
    txt(s,0.95,y+0.82,11.6,0.5,(st["text"],10,INK),sp=1.05)
txt(s,0.7,7.0,12,0.3,("Human (steel) → Agent (coral): the boundary slides right as the level rises — the human holds the gate until L2.",11,MUT,True))

# 16 · GOALS
s=slide(); head(s,"The Dial","Everything you asked for — your goals, on the gentler path",num="15")
txt(s,0.7,1.5,12.2,0.6,("12 of your 16 Data & Integration goals and 3 of 4 DevOps objectives are fully met at Level 0 — nothing agentic. Nothing needs Level 2. The dial governs pace, never coverage.",13,INK,True),sp=1.15)
allg=[("Data",g) for g in D["goalsB21"]]+[("Integration",g) for g in D["goalsB22"]]+[("DevOps",g) for g in D["goalsC2"]]
rows=[["Ref","Your stated goal","Met at","Horizon"]]
for pref,g in allg:
    rows.append([g["ref"], g["goal"][:60], ("Level 1" if g["metAt"]==1 else "Level 0"), g["horizon"]])
table(s,0.7,2.3,12.0,rows,[1.5,6.3,1.3,2.9],rowh=0.3,body_size=8.5,head_size=9.5)

# 17 · JOURNEY
s=slide(); head(s,"The Journey","Re-anchored to the 30 Oct award — both lanes, four checkpoints",num="16")
hz=D["horizons"]; hb={h["id"]:h for h in D["horizonBoard"]}
txt(s,3.4,1.5,4.5,0.3,("LANE 1 · RUN & DELIVER",10.5,TEC,True)); txt(s,8.4,1.5,4.5,0.3,("LANE 2 · IMPROVE & EVOLVE",10.5,ACC,True))
for i,h in enumerate(hz):
    y=1.85+i*1.28; row=hb[h["id"]]
    box(s,0.7,y,2.45,1.16); txt(s,0.85,y+0.14,2.2,0.9,[(f"{h['months']} · {h['date']}",11.5,INK,True),(h["stage"],10,MUT,False)],sp=1.05)
    box(s,3.35,y,4.8,1.16,fill=RGBColor(0xED,0xF1,0xF6),line=RGBColor(0xB6,0xC5,0xD5)); txt(s,3.5,y+0.1,4.5,1.0,"• "+"\n• ".join(row["lane1"][:3]),9,INK,sp=1.05)
    box(s,8.35,y,4.28,1.16,fill=RGBColor(0xF7,0xE9,0xE9),line=RGBColor(0xE4,0xB9,0xBD)); txt(s,8.5,y+0.1,3.95,1.0,"• "+"\n• ".join(row["lane2"][:3]),9,INK,sp=1.05)

# 18 · COST OF CAUTION
s=slide(); head(s,"The cost of caution","Shown, not hidden — what each target needs",num="17")
rows=[["RFP target","L0","L1","L2"]]+[list(r) for r in D["cautionRows"]]
table(s,0.7,1.6,12.0,rows,[5.2,2.4,2.4,2.0],rowh=0.42,body_size=10.5)
foot(s,"Three headline targets — sub-30-min workspaces, minutes-onboarding, top-quartile DORA — are Level-1 outcomes. The dial is yours.")

# 19 · PROOF
s=slide(); head(s,"Why this is low-risk","Proven takeover · fabric already built · goals covered",num="18")
for i,pc in enumerate(D["proofCards"]):
    x=0.7+(i%2)*6.15; y=1.7+(i//2)*2.5
    box(s,x,y,5.95,2.3)
    chip(s,x+0.2,y+0.2,pc["chip"],GRN)
    txt(s,x+0.2,y+0.6,5.6,0.5,(pc["stat"],22,ACC,True))
    txt(s,x+0.2,y+1.25,5.6,0.95,(pc["body"],10,MUT),sp=1.1)

# 20 · YOUR PART + START
s=slide(); head(s,"What we need from Pandora","Deliberately lighter — and how we start",num="19")
txt(s,0.7,1.55,12,0.3,("WHAT WE NEED FROM PANDORA",11,MUT,True))
bullets(s,0.7,1.9,12,2.0,D["pandoraAsks"],size=12,gap=1.35,mcolor=GRN)
txt(s,0.7,4.15,12,0.3,("HOW WE START — FROM AWARD TO THE FIRST DIAL REVIEW",11,MUT,True))
for i,sp_ in enumerate(D["startPlan"]):
    y=4.5+i*0.7
    box(s,0.7,y,12,0.62); txt(s,0.9,y+0.14,3.2,0.4,(sp_["when"],11,ACC,True)); txt(s,4.2,y+0.11,8.3,0.5,(sp_["what"],10.5,MUT),sp=1.05)

# 21 · THE ASK
s=slide(bg=INK)
txt(s,0.9,0.8,11,0.35,("20 · THE ASK",12,ACC,True))
txt(s,0.9,1.2,11.4,0.7,("One decision starts the gentler path",30,WHITE,True))
for i,(t,b) in enumerate(D["askCards"]):
    y=2.3+i*1.1
    c=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.9),Inches(y),Inches(11.5),Inches(0.98)); c.fill.background(); c.line.color.rgb=RGBColor(0x3A,0x3E,0x45); c.line.width=Pt(0.75); c.shadow.inherit=False; c.adjustments[0]=0.06
    txt(s,1.2,y+0.14,2.7,0.4,(t.upper(),11.5,ACC,True))
    txt(s,4.0,y+0.12,8.1,0.7,(b,12,RGBColor(0xC9,0xCB,0xD0)),sp=1.1)
txt(s,0.9,y+1.2,11.5,0.8,(D["closeLine"],15,RGBColor(0xE8,0xE6,0xE1),True),sp=1.15)

prs.save(OUT)
print(f"saved {len(prs.slides.__iter__.__self__._sldIdLst)} slides -> {OUT}")
